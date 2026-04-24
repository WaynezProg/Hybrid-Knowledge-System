"""pptx parser."""

from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError

from hks.errors import ExitCode, KSError
from hks.ingest.fingerprint import ParserFlags
from hks.ingest.models import ParsedDocument
from hks.ingest.office_common import (
    Segment,
    SkippedSegment,
    build_placeholder,
    merge_skipped,
    to_markdown_table,
)


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        text = "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text.strip()
        )
        return text.strip()
    return ""


def _append_placeholder(
    segments: list[Segment],
    skipped_segments: list[SkippedSegment],
    kind: str,
    payload: str | None = None,
) -> None:
    merge_skipped(skipped_segments, SkippedSegment(type=kind))  # type: ignore[arg-type]
    segments.append(Segment(kind="placeholder", text=build_placeholder(kind, payload)))  # type: ignore[arg-type]


def _append_table(shape: Any, *, slide_index: int, segments: list[Segment]) -> None:
    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
    if not rows:
        return
    header = [cell or f"col_{index + 1}" for index, cell in enumerate(rows[0])]
    body_rows = [row for row in rows[1:] if any(cell for cell in row)]
    if not body_rows:
        segments.append(
            Segment(
                kind="table_row",
                text=to_markdown_table(header, []),
                metadata={"slide_index": slide_index},
            )
        )
        return
    for row in body_rows:
        segments.append(
            Segment(
                kind="table_row",
                text=to_markdown_table(header, [row]),
                metadata={"slide_index": slide_index},
            )
        )


def _notes_text(slide: Any) -> str:
    texts: list[str] = []
    for placeholder in slide.notes_slide.placeholders:
        if placeholder.placeholder_format.idx != 2:
            continue
        if getattr(placeholder, "has_text_frame", False):
            text = placeholder.text.strip()
            if text:
                texts.append(text)
    return "\n".join(texts).strip()


def parse(path: Path, flags: ParserFlags) -> ParsedDocument:
    try:
        presentation = Presentation(str(path))
    except (PackageNotFoundError, ValueError, zipfile.BadZipFile) as exc:
        raise KSError(
            f"corrupt file: {path}",
            exit_code=ExitCode.DATAERR,
            code="CORRUPT",
            details=[str(exc)],
        ) from exc

    segments: list[Segment] = []
    skipped_segments: list[SkippedSegment] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_segments_before = len(segments)
        segments.append(
            Segment(
                kind="slide_header",
                text=f"## Slide {slide_index}",
                metadata={"slide_index": slide_index},
            )
        )

        title_shape = slide.shapes.title
        title_text = title_shape.text.strip() if title_shape is not None else ""
        if title_text:
            segments.append(
                Segment(
                    kind="heading",
                    text=f"### {title_text}",
                    metadata={"level": 3, "slide_index": slide_index},
                )
            )

        text_segments: list[Segment] = []
        table_segments: list[Segment] = []
        placeholder_segments: list[Segment] = []
        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue
            if getattr(shape, "has_table", False):
                _append_table(shape, slide_index=slide_index, segments=table_segments)
                continue

            shape_text = _shape_text(shape)
            if shape_text:
                text_segments.append(
                    Segment(
                        kind="paragraph",
                        text=shape_text,
                        metadata={"slide_index": slide_index},
                    )
                )
                continue

            shape_type = shape.shape_type
            if shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}:
                merge_skipped(skipped_segments, SkippedSegment(type="image"))
                placeholder_segments.append(
                    Segment(
                        kind="placeholder",
                        text=build_placeholder("image", shape.name.strip()),
                        metadata={"slide_index": slide_index},
                    )
                )
            elif shape_type in {
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            }:
                merge_skipped(skipped_segments, SkippedSegment(type="embedded_object"))
                placeholder_segments.append(
                    Segment(
                        kind="placeholder",
                        text=build_placeholder("embedded_object", shape.name.strip()),
                        metadata={"slide_index": slide_index},
                    )
                )
            elif shape_type in {MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.IGX_GRAPHIC}:
                merge_skipped(skipped_segments, SkippedSegment(type="smartart"))
                placeholder_segments.append(
                    Segment(
                        kind="placeholder",
                        text=build_placeholder("smartart", ""),
                        metadata={"slide_index": slide_index},
                    )
                )
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                merge_skipped(skipped_segments, SkippedSegment(type="chart"))
                placeholder_segments.append(
                    Segment(
                        kind="placeholder",
                        text=build_placeholder("chart"),
                        metadata={"slide_index": slide_index},
                    )
                )
            elif shape_type in {MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO}:
                media_kind = (
                    "audio"
                    if "audio" in shape.name.lower() or "sound" in shape.name.lower()
                    else "video"
                )
                merge_skipped(skipped_segments, SkippedSegment(type=media_kind))  # type: ignore[arg-type]
                placeholder_segments.append(
                    Segment(
                        kind="placeholder",
                        text=build_placeholder(media_kind),  # type: ignore[arg-type]
                        metadata={"slide_index": slide_index},
                    )
                )

        segments.extend(text_segments)
        segments.extend(table_segments)
        segments.extend(placeholder_segments)

        if flags.pptx_notes:
            notes_text = _notes_text(slide)
            if notes_text:
                segments.append(
                    Segment(
                        kind="notes",
                        text=notes_text,
                        metadata={"slide_index": slide_index},
                    )
                )

        if len(segments) == slide_segments_before + 1:
            merge_skipped(skipped_segments, SkippedSegment(type="empty_slide"))

    try:
        with zipfile.ZipFile(path) as archive:
            if any(name.endswith("vbaProject.bin") for name in archive.namelist()):
                merge_skipped(skipped_segments, SkippedSegment(type="macros"))
                segments.append(Segment(kind="placeholder", text=build_placeholder("macros")))
    except (OSError, zipfile.BadZipFile) as exc:
        raise KSError(
            f"corrupt file: {path}",
            exit_code=ExitCode.DATAERR,
            code="CORRUPT",
            details=[str(exc)],
        ) from exc

    return ParsedDocument(
        title=path.stem.replace("-", " ").replace("_", " ").strip() or path.stem,
        body="",
        format="pptx",
        segments=segments,
        skipped_segments=skipped_segments,
    )
