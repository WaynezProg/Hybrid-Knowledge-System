"""Programmatic Office fixture builder for Phase 2 tests.

Run via `make fixtures` or `python tests/fixtures/build_office.py`.
Generates deterministic Office fixtures under:

- `tests/fixtures/valid/docx/`
- `tests/fixtures/valid/xlsx/`
- `tests/fixtures/valid/pptx/`
- `tests/fixtures/broken/office/`
"""

from __future__ import annotations

import base64
import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt

FIXTURES_ROOT = Path(__file__).resolve().parent
VALID_ROOT = FIXTURES_ROOT / "valid"
BROKEN_ROOT = FIXTURES_ROOT / "broken" / "office"
DOCX_ROOT = VALID_ROOT / "docx"
XLSX_ROOT = VALID_ROOT / "xlsx"
PPTX_ROOT = VALID_ROOT / "pptx"
ASSET_ROOT = FIXTURES_ROOT / "_assets"

PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO8B2i0AAAAASUVORK5CYII="
)
OLE_MAGIC = bytes.fromhex("D0CF11E0A1B11AE1")
SHEET_NS = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}


def ensure_dirs() -> None:
    for path in (DOCX_ROOT, XLSX_ROOT, PPTX_ROOT, BROKEN_ROOT, ASSET_ROOT):
        path.mkdir(parents=True, exist_ok=True)


def _asset_png() -> Path:
    image_path = ASSET_ROOT / "tiny.png"
    image_path.write_bytes(PNG_BYTES)
    return image_path


def _replace_zip_entries(path: Path, updates: dict[str, bytes]) -> None:
    temp_path = path.with_suffix(f"{path.suffix}.tmp")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(temp_path, "w") as target:
        for info in source.infolist():
            if info.filename in updates:
                continue
            target.writestr(info, source.read(info.filename))
        for filename, payload in updates.items():
            target.writestr(filename, payload)
    temp_path.replace(path)


def _inject_macro(path: Path, *, package_subdir: str) -> None:
    _replace_zip_entries(path, {f"{package_subdir}/vbaProject.bin": b"fake-vba-project"})


def _write_formula_cache(path: Path, *, cached_values: dict[str, str]) -> None:
    with zipfile.ZipFile(path) as archive:
        sheet_xml = archive.read("xl/worksheets/sheet1.xml")
    root = ET.fromstring(sheet_xml)
    changed = False
    for cell in root.findall(".//main:c", SHEET_NS):
        ref = cell.attrib.get("r", "")
        if ref not in cached_values:
            continue
        value = cell.find("main:v", SHEET_NS)
        if value is None:
            value = ET.SubElement(cell, f"{{{SHEET_NS['main']}}}v")
        value.text = cached_values[ref]
        changed = True
    if changed:
        _replace_zip_entries(
            path,
            {"xl/worksheets/sheet1.xml": ET.tostring(root, encoding="utf-8")},
        )


def _make_oversized_office_file(source: Path, destination: Path, *, payload_mb: int = 2) -> None:
    shutil.copy2(source, destination)
    padding = b"x" * (payload_mb * 1024 * 1024)
    _replace_zip_entries(destination, {"xl/padding.bin": padding})


def _set_docx_alt_text(paragraph, description: str) -> None:
    run = paragraph.add_run()
    inline_shape = run.add_picture(str(_asset_png()), width=Inches(1))
    inline_shape._inline.docPr.set("descr", description)


def _set_notes(slide, text: str) -> None:
    for placeholder in slide.notes_slide.placeholders:
        if placeholder.placeholder_format.idx == 2:
            placeholder.text = text
            return
    raise RuntimeError("pptx notes placeholder idx=2 missing")


def build_docx() -> None:
    plain = Document()
    plain.add_heading("Atlas Office Summary", level=1)
    plain.add_paragraph("Atlas roadmap is on track with two unresolved procurement risks.")
    plain.add_paragraph("Key dependency is the mobile gateway migration scheduled for June.")
    plain.add_paragraph("Escalate vendor contract by Friday.", style="List Bullet")
    plain.add_paragraph("Confirm rollout checklist with finance.", style="List Bullet")
    plain.save(DOCX_ROOT / "plain.docx")

    with_table = Document()
    with_table.add_heading("Delivery Status", level=1)
    with_table.add_paragraph("Current milestone table:")
    table = with_table.add_table(rows=3, cols=3)
    table.rows[0].cells[0].text = "team"
    table.rows[0].cells[1].text = "status"
    table.rows[0].cells[2].text = "owner"
    table.rows[1].cells[0].text = "search"
    table.rows[1].cells[1].text = "green"
    table.rows[1].cells[2].text = "Iris"
    table.rows[2].cells[0].text = "billing"
    table.rows[2].cells[1].text = "yellow"
    table.rows[2].cells[2].text = "Noah"
    with_table.save(DOCX_ROOT / "with_table.docx")

    with_image = Document()
    with_image.add_heading("Field Visit", level=1)
    with_image.add_paragraph("Photo evidence captured during the warehouse walkthrough.")
    image_paragraph = with_image.add_paragraph()
    _set_docx_alt_text(image_paragraph, "Figure 1: sample")
    with_image.save(DOCX_ROOT / "with_image.docx")


def build_xlsx() -> None:
    single = Workbook()
    sheet = single.active
    sheet.title = "Summary"
    sheet.append(["id", "project", "risk"])
    for index in range(1, 11):
        sheet.append([index, f"Atlas-{index}", f"risk-{index}"])
    single.save(XLSX_ROOT / "single_sheet.xlsx")

    multi = Workbook()
    summary = multi.active
    summary.title = "Summary"
    summary.append(["project", "owner", "status"])
    summary.append(["Atlas", "Iris", "green"])
    budget = multi.create_sheet("Budget")
    budget.append(["quarter", "amount"])
    budget.append(["Q1", 120])
    budget.append(["Q2", 135])
    risks = multi.create_sheet("Risks")
    risks.append(["severity", "issue"])
    risks.append(["high", "vendor delay"])
    risks.append(["medium", "QA staffing"])
    multi.save(XLSX_ROOT / "multi_sheet.xlsx")

    with_formula = Workbook()
    formula_sheet = with_formula.active
    formula_sheet.title = "Metrics"
    formula_sheet.append(["base", "delta", "total"])
    formula_sheet.append([2, 3, "=SUM(A2:B2)"])
    formula_sheet.append([5, 8, "=SUM(A3:B3)"])
    chart = BarChart()
    chart.title = "Quarterly Totals"
    chart.add_data(Reference(formula_sheet, min_col=3, min_row=1, max_row=3), titles_from_data=True)
    formula_sheet.add_chart(chart, "E2")
    with_formula.save(XLSX_ROOT / "with_formula.xlsx")
    _write_formula_cache(XLSX_ROOT / "with_formula.xlsx", cached_values={"C2": "5"})
    _inject_macro(XLSX_ROOT / "with_formula.xlsx", package_subdir="xl")


def build_pptx() -> None:
    plain = Presentation()
    for index, (title, body) in enumerate(
        [
            ("Atlas Overview", "Atlas enters pilot launch next month."),
            ("Milestones", "Search indexing and billing validation are complete."),
            ("Risks", "Largest risk remains vendor lead time."),
            ("Dependencies", "Mobile gateway migration gates downstream release."),
            ("Next Step", "Approve rollout checklist and train support."),
        ],
        start=1,
    ):
        slide = plain.slides.add_slide(plain.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Slide {index}: {body}"
    plain.save(PPTX_ROOT / "plain.pptx")

    with_notes = Presentation()
    for title, body, notes in [
        ("Kickoff", "Launch meeting on Monday.", "Speaker notes mention Monday owner = Iris."),
        (
            "Budget",
            "Budget review follows procurement sign-off.",
            "Speaker notes mention budget delta 15.",
        ),
        ("Risks", "Vendor lead time still volatile.", "Speaker notes mention fallback supplier."),
    ]:
        slide = with_notes.slides.add_slide(with_notes.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        _set_notes(slide, notes)
    with_notes.save(PPTX_ROOT / "with_notes.pptx")

    with_table_image = Presentation()
    slide = with_table_image.slides.add_slide(with_table_image.slide_layouts[5])
    slide.shapes.title.text = "Capacity Table"
    table_shape = slide.shapes.add_table(
        3,
        2,
        PptxInches(1),
        PptxInches(1.8),
        PptxInches(5),
        PptxInches(1.5),
    )
    table = table_shape.table
    table.cell(0, 0).text = "team"
    table.cell(0, 1).text = "capacity"
    table.cell(1, 0).text = "search"
    table.cell(1, 1).text = "high"
    table.cell(2, 0).text = "billing"
    table.cell(2, 1).text = "medium"

    image_slide = with_table_image.slides.add_slide(with_table_image.slide_layouts[6])
    image_slide.shapes.add_picture(
        str(_asset_png()),
        PptxInches(1),
        PptxInches(1),
        width=PptxInches(2),
    )

    mixed_slide = with_table_image.slides.add_slide(with_table_image.slide_layouts[1])
    mixed_slide.shapes.title.text = "Floor Plan"
    body = mixed_slide.placeholders[1]
    body.text = "Photo attached for reference."
    body.text_frame.paragraphs[0].font.size = Pt(20)
    mixed_slide.shapes.add_picture(
        str(_asset_png()),
        PptxInches(5.5),
        PptxInches(1.5),
        width=PptxInches(1.5),
    )
    with_table_image.save(PPTX_ROOT / "with_table_image.pptx")


def build_broken() -> None:
    encrypted = BROKEN_ROOT / "encrypted.pptx"
    encrypted.write_bytes(OLE_MAGIC + (b"\x00" * 4096))

    corrupt = BROKEN_ROOT / "corrupt.xlsx"
    shutil.copy2(XLSX_ROOT / "single_sheet.xlsx", corrupt)
    corrupt.write_bytes(corrupt.read_bytes()[:128])

    (BROKEN_ROOT / "empty.docx").write_bytes(b"")

    timeout_bomb = Document()
    timeout_bomb.add_heading("Timeout Bomb", level=1)
    for index in range(50000):
        timeout_bomb.add_paragraph(f"timeout paragraph {index}: " + ("payload " * 20))
    timeout_bomb.save(BROKEN_ROOT / "timeout_bomb.docx")

    oversized = BROKEN_ROOT / "oversized.xlsx"
    _make_oversized_office_file(XLSX_ROOT / "single_sheet.xlsx", oversized, payload_mb=2)


def build_all() -> None:
    ensure_dirs()
    build_docx()
    build_xlsx()
    build_pptx()
    build_broken()


if __name__ == "__main__":
    build_all()
